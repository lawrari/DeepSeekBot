import io
from docx import Document
import pandas as pd
import openpyxl
from openpyxl.drawing.image import Image as XlsxImage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class FilesToText:

    @staticmethod
    def extract_docx_content(docx_bytes):
        """
        Extract text and images from a DOCX file with comprehensive text extraction.
        
        Args:
            docx_bytes (bytes): The content of the DOCX file as bytes
            
        Returns:
            dict: A dictionary with 'text' containing the full document text and
                'images' containing a list of image bytes
        """
        # Create a BytesIO object from the docx bytes
        docx_file = io.BytesIO(docx_bytes.read())
        
        # Load the document
        doc = Document(docx_file)
        
        # Extract text from paragraphs
        paragraph_text = []
        for paragraph in doc.paragraphs:
            paragraph_text.append(paragraph.text)
        
        # Extract text from tables
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    # Get text from each paragraph in the cell
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            row_text.append(paragraph.text)
                if row_text:
                    table_text.append(" | ".join(row_text))
        
        # Extract text from headers
        header_text = []
        for section in doc.sections:
            header = section.header
            for paragraph in header.paragraphs:
                if paragraph.text.strip():
                    header_text.append(paragraph.text)
        
        # Extract text from footers
        footer_text = []
        for section in doc.sections:
            footer = section.footer
            for paragraph in footer.paragraphs:
                if paragraph.text.strip():
                    footer_text.append(paragraph.text)
        
        # Combine all text sections with appropriate separators
        all_text_parts = []
        
        if paragraph_text:
            all_text_parts.append("\n".join(paragraph_text))
        
        if table_text:
            all_text_parts.append("\n=== TABLES ===\n" + "\n".join(table_text))
        
        if header_text:
            all_text_parts.append("\n=== HEADERS ===\n" + "\n".join(header_text))
        
        if footer_text:
            all_text_parts.append("\n=== FOOTERS ===\n" + "\n".join(footer_text))
        
        # Extract images
        image_bytes_list = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_part = rel.target_part
                if hasattr(image_part, 'blob'):
                    # Get the image bytes
                    img_bytes = image_part.blob
                    image_bytes_list.append(img_bytes)
        
        # Create the result dictionary
        result = {
            "text": "\n\n".join(all_text_parts),
            "images": image_bytes_list
        }
        
        return result

    @staticmethod
    def extract_xlsx_content(xlsx_bytes):
        """
        Extract text and images from an XLSX file.
        
        Args:
            xlsx_bytes (bytes): The content of the XLSX file as bytes
            
        Returns:
            dict: A dictionary with 'text' containing the sheet data as text and
                'images' containing a list of image bytes
        """
        # Create BytesIO objects
        xlsx_file = io.BytesIO(xlsx_bytes.read())
        
        # Extract text using pandas (for all sheets)
        text_output = []
        xlsx_file.seek(0)
        xl = pd.ExcelFile(xlsx_file)
        sheet_names = xl.sheet_names
        
        for sheet_name in sheet_names:
            df = pd.read_excel(xlsx_file, sheet_name=sheet_name)
            text_output.append(f"=== SHEET: {sheet_name} ===")
            text_output.append(df.to_string(index=False))
        
        # Extract images using openpyxl
        image_bytes_list = []
        xlsx_file.seek(0)
        workbook = openpyxl.load_workbook(xlsx_file)
        
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            
            # Extract images from the worksheet
            for image in worksheet._images:
                # Get the image data
                img_data = image._data()
                if hasattr(image, '_data'):
                    image_bytes_list.append(img_data)
        
        # For embedded charts (which might contain images)
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            if hasattr(worksheet, '_charts'):
                for chart in worksheet._charts:
                    if hasattr(chart, 'images'):
                        for img in chart.images:
                            if hasattr(img, '_data'):
                                image_bytes_list.append(img._data())
        
        # Extract any other drawing objects that might be images
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            if hasattr(worksheet, 'drawings'):
                for drawing in worksheet.drawings.values():
                    if isinstance(drawing, XlsxImage) and hasattr(drawing, 'ref'):
                        try:
                            img_data = drawing._data()
                            image_bytes_list.append(img_data)
                        except:
                            pass
        
        # Create the result dictionary
        result = {
            "text": "\n\n".join(text_output),
            "images": image_bytes_list
        }
        
        return result

    @staticmethod
    def extract_pptx_content(pptx_bytes):
        """
        Extract text and images from a PPTX file.
        
        Args:
            pptx_bytes (bytes): The content of the PPTX file as bytes
            
        Returns:
            dict: A dictionary with 'text' containing the slide text and
                'images' containing a list of image bytes
        """
        # Create a BytesIO object from the pptx bytes
        pptx_file = io.BytesIO(pptx_bytes.read())
        
        # Load the presentation
        presentation = Presentation(pptx_file)
        
        # Extract text
        text_output = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            slide_text.append(f"=== SLIDE {slide_num} ===")
            
            # Extract slide title if exists
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"Title: {slide.shapes.title.text}")
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and shape != slide.shapes.title:  # Skip title we already added
                        slide_text.append(text)
                
                # Extract text from tables
                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        row_cells = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_cells.append(cell.text.strip())
                        if row_cells:
                            table_rows.append(" | ".join(row_cells))
                    if table_rows:
                        slide_text.append("Table:")
                        slide_text.extend(table_rows)
                        
                # Extract text from grouped shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for subshape in shape.shapes:
                        if hasattr(subshape, "text") and subshape.text.strip():
                            slide_text.append(subshape.text.strip())
            
            # Add notes if they exist
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                slide_text.append(f"Notes: {slide.notes_slide.notes_text_frame.text.strip()}")
                
            if len(slide_text) > 1:  # If we have more than just the slide title
                text_output.extend(slide_text)
        
        # Extract images
        image_bytes_list = []
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                # Check if shape is a picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    # Get image bytes in their original format
                    image_bytes = image.blob
                    image_bytes_list.append(image_bytes)
                
                # Check grouped shapes for pictures
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for subshape in shape.shapes:
                        if hasattr(subshape, "shape_type") and subshape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image = subshape.image
                            image_bytes = image.blob
                            image_bytes_list.append(image_bytes)
        
        # Create the result dictionary
        result = {
            "text": "\n".join(text_output),
            "images": image_bytes_list
        }
        
        return result
    
    def file_to_text(file_bytes: bytes, encoding: str = "utf-8") -> str:
        """
        Converts bytes of a text-based file into a string.

        Args:
            file_bytes (bytes): The file content in bytes.
            encoding (str, optional): The text encoding (default is 'utf-8').

        Returns:
            str: The decoded text content.
        """
        try:
            return file_bytes.read().decode(encoding)
        except UnicodeDecodeError:
            return file_bytes.decode("utf-8", errors="replace")
